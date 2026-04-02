"""
生成PPT展示部门树状图
"""
import sqlite3
import os
import shutil
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def get_departments_from_db():
    """从数据库获取所有部门数据"""
    conn = sqlite3.connect('employee.db')
    cursor = conn.cursor()
    cursor.execute('SELECT id, parent_id, name, level FROM department ORDER BY level, id')
    departments = cursor.fetchall()
    conn.close()
    return departments


def build_department_tree(departments):
    """构建部门树结构"""
    # 创建节点映射
    nodes = {d[0]: {'id': d[0], 'parent_id': d[1], 'name': d[2], 'level': d[3], 'children': []}
             for d in departments}

    # 构建树
    roots = []
    for d in departments:
        node = nodes[d[0]]
        if d[1] is None:
            roots.append(node)
        else:
            parent = nodes.get(d[1])
            if parent:
                parent['children'].append(node)

    return roots


def count_tree_nodes(node):
    """计算树的总节点数"""
    count = 1
    for child in node['children']:
        count += count_tree_nodes(child)
    return count


def get_max_width_at_level(node, level=0, widths=None):
    """计算每层的最大宽度"""
    if widths is None:
        widths = {}
    widths[level] = widths.get(level, 0) + 1
    for child in node['children']:
        get_max_width_at_level(child, level + 1, widths)
    return widths


def get_level_widths(nodes):
    """计算每层的最大宽度（同层部门数量）"""
    widths = {}

    def traverse(node, level):
        widths[level] = widths.get(level, 0) + 1
        for child in node['children']:
            traverse(child, level + 1)

    for root in nodes:
        traverse(root, 0)
    return widths


def draw_tree(presentation, nodes):
    """在PPT上绘制树状图"""
    # 添加空白幻灯片
    blank_layout = presentation.slide_layouts[6]  # 空白布局
    slide = presentation.slides.add_slide(blank_layout)

    # 配置参数 - 部门树只占幻灯片高度的40%
    slide_height = Inches(7.5)
    slide_width = Inches(13.333)  # 幻灯片宽度 (16:9)
    tree_height = slide_height * 0.4  # 40%高度 = 3英寸
    box_height = Inches(0.35)
    margin = Inches(0.3)

    # 计算每层部门数量
    level_widths = get_level_widths(nodes)
    max_level = max(level_widths.keys()) if level_widths else 0

    # 根据最宽层计算最大可用节点框宽度
    available_width = slide_width - 2 * margin
    max_count = max(level_widths.values()) if level_widths else 1
    max_box_width = (available_width - max_count * Inches(0.15)) / max_count  # 留0.15英寸间距
    max_box_width = min(max_box_width, Inches(2.0))  # 最大宽度不超过2英寸

    # 层级间距
    level_height = tree_height / (max_level + 1)
    start_y = Inches(0.2)

    # 颜色配置（按层级）
    level_colors = [
        RGBColor(0x4A, 0x90, 0xD9),  # 一层 - 蓝色
        RGBColor(0x5C, 0xB8, 0x5C),  # 二层 - 绿色
        RGBColor(0xF0, 0xAD, 0x4E),  # 三层 - 橙色
        RGBColor(0xD9, 0x5F, 0x5F),  # 四层 - 红色
        RGBColor(0x9B, 0x59, 0xB6),  # 五层 - 紫色
    ]

    # 存储所有节点的位置信息
    node_positions = {}
    box_widths = {}  # 存储每个节点的宽度

    def calculate_positions(node, depth, x_start, x_end):
        """递归计算每个节点的位置"""
        y = start_y + depth * level_height

        # 计算当前节点的节点框宽度
        # 宽度 = 分配空间 - 间距，但不超过最大宽度
        space = x_end - x_start
        gap = Inches(0.1)  # 节点间距
        box_w = min(space - gap, max_box_width)
        box_w = max(box_w, Inches(0.6))  # 最小宽度0.6英寸

        # 当前节点的x位置在分配范围的中间
        x_center = (x_start + x_end) / 2
        x = x_center - box_w / 2

        node_positions[node['id']] = {
            'x': x,
            'y': y,
            'x_center': x_center,
            'y_center': y + box_height / 2,
            'box_right': x + box_w,
            'box_bottom': y + box_height,
            'name': node['name'],
            'level': depth
        }
        box_widths[node['id']] = box_w

        # 递归处理子节点
        if node['children']:
            num_children = len(node['children'])
            child_width = (x_end - x_start) / num_children
            for i, child in enumerate(node['children']):
                child_x_start = x_start + i * child_width
                child_x_end = child_x_start + child_width
                calculate_positions(child, depth + 1, child_x_start, child_x_end)

    for root in nodes:
        calculate_positions(root, 0, margin, slide_width - margin)

    # 绘制连接线和节点
    def draw_tree_recursive(node):
        """递归绘制连接线"""
        pos = node_positions[node['id']]
        for child in node['children']:
            child_pos = node_positions[child['id']]
            draw_line(slide, pos['x_center'], pos['box_bottom'],
                     child_pos['x_center'], child_pos['y'])
            draw_tree_recursive(child)

    for root in nodes:
        draw_tree_recursive(root)

    # 绘制所有节点
    for node_id, pos in node_positions.items():
        draw_node(slide, pos['x'], pos['y'], box_widths[node_id], box_height,
                 pos['name'], level_colors[pos['level']])


def draw_tree_with_return(presentation, nodes):
    """在PPT上绘制树状图，并返回节点位置信息用于生成文档"""
    # 添加空白幻灯片
    blank_layout = presentation.slide_layouts[6]  # 空白布局
    slide = presentation.slides.add_slide(blank_layout)

    # 配置参数 - 部门树只占幻灯片高度的40%
    slide_height = Inches(7.5)
    slide_width = Inches(13.333)  # 幻灯片宽度 (16:9)
    tree_height = slide_height * 0.4  # 40%高度 = 3英寸
    box_height = Inches(0.35)
    margin = Inches(0.3)

    # 计算每层部门数量
    level_widths = get_level_widths(nodes)
    max_level = max(level_widths.keys()) if level_widths else 0

    # 根据最宽层计算最大可用节点框宽度
    available_width = slide_width - 2 * margin
    max_count = max(level_widths.values()) if level_widths else 1
    max_box_width = (available_width - max_count * Inches(0.15)) / max_count  # 留0.15英寸间距
    max_box_width = min(max_box_width, Inches(2.0))  # 最大宽度不超过2英寸

    # 层级间距
    level_height = tree_height / (max_level + 1)
    start_y = Inches(0.2)

    # 颜色配置（按层级）
    level_colors = [
        RGBColor(0x4A, 0x90, 0xD9),  # 一层 - 蓝色
        RGBColor(0x5C, 0xB8, 0x5C),  # 二层 - 绿色
        RGBColor(0xF0, 0xAD, 0x4E),  # 三层 - 橙色
        RGBColor(0xD9, 0x5F, 0x5F),  # 四层 - 红色
        RGBColor(0x9B, 0x59, 0xB6),  # 五层 - 紫色
    ]

    # 存储所有节点的位置信息
    node_positions = {}
    box_widths = {}  # 存储每个节点的宽度

    def calculate_positions(node, depth, x_start, x_end):
        """递归计算每个节点的位置"""
        y = start_y + depth * level_height

        # 计算当前节点的节点框宽度
        # 宽度 = 分配空间 - 间距，但不超过最大宽度
        space = x_end - x_start
        gap = Inches(0.1)  # 节点间距
        box_w = min(space - gap, max_box_width)
        box_w = max(box_w, Inches(0.6))  # 最小宽度0.6英寸

        # 当前节点的x位置在分配范围的中间
        x_center = (x_start + x_end) / 2
        x = x_center - box_w / 2

        node_positions[node['id']] = {
            'x': x,
            'y': y,
            'x_center': x_center,
            'y_center': y + box_height / 2,
            'box_right': x + box_w,
            'box_bottom': y + box_height,
            'name': node['name'],
            'level': depth,
            'box_width': box_w
        }
        box_widths[node['id']] = box_w

        # 递归处理子节点
        if node['children']:
            num_children = len(node['children'])
            child_width = (x_end - x_start) / num_children
            for i, child in enumerate(node['children']):
                child_x_start = x_start + i * child_width
                child_x_end = child_x_start + child_width
                calculate_positions(child, depth + 1, child_x_start, child_x_end)

    for root in nodes:
        calculate_positions(root, 0, margin, slide_width - margin)

    # 绘制连接线和节点
    def draw_tree_recursive(node):
        """递归绘制连接线"""
        pos = node_positions[node['id']]
        for child in node['children']:
            child_pos = node_positions[child['id']]
            draw_line(slide, pos['x_center'], pos['box_bottom'],
                     child_pos['x_center'], child_pos['y'])
            draw_tree_recursive(child)

    for root in nodes:
        draw_tree_recursive(root)

    # 绘制所有节点
    for node_id, pos in node_positions.items():
        draw_node(slide, pos['x'], pos['y'], box_widths[node_id], box_height,
                 pos['name'], level_colors[pos['level']])

    return node_positions, level_colors, slide, level_widths


def draw_node(slide, x, y, width, height, text, fill_color):
    """绘制一个节点"""
    shape = slide.shapes.add_shape(1, x, y, width, height)  # 1 = 矩形
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # 处理文字：超过10个字用省略号
    if len(text) > 10:
        display_text = text[:10] + '...'
    else:
        display_text = text

    # 根据节点框宽度和文字数量计算合适的字体大小
    # 宽度单位是EMU，转换为英寸
    width_inches = width / 914400  # 1英寸 = 914400 EMU

    # 计算字体大小：根据宽度和文字长度动态调整
    # 基础字体大小11pt，每个字符约占0.1英寸（11pt字体）
    char_count = len(display_text)
    # 可用宽度减去边距（每边约0.05英寸）
    available_width = width_inches - 0.1

    # 计算合适的字体大小
    if char_count > 0:
        # 每个字符的最大宽度（英寸）
        max_char_width = available_width / char_count
        # 字体大小（pt）与字符宽度的近似关系：1pt ≈ 0.014英寸
        font_size = min(11, max(7, int(max_char_width / 0.014)))
    else:
        font_size = 11

    # 添加文字
    tf = shape.text_frame
    tf.word_wrap = False  # 禁用自动换行，确保单行显示
    p = tf.paragraphs[0]
    p.text = display_text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(6)


def draw_line(slide, x1, y1, x2, y2):
    """绘制连接线 - 使用Freeform形状"""
    from pptx.oxml.ns import nsmap
    from pptx.oxml import parse_xml

    # 使用最简单的方式：绘制一个很细的矩形作为线条
    # 对于垂直线
    if abs(x2 - x1) < Inches(0.1):
        # 垂直线
        left = x1 - Inches(0.01)
        top = min(y1, y2)
        width = Inches(0.02)
        height = abs(y2 - y1)
    elif abs(y2 - y1) < Inches(0.1):
        # 水平线
        left = min(x1, x2)
        top = y1 - Inches(0.01)
        width = abs(x2 - x1)
        height = Inches(0.02)
    else:
        # 斜线：用折线方式
        draw_polyline(slide, x1, y1, x2, y2)
        return

    line = slide.shapes.add_shape(1, left, top, width, height)  # 1 = 矩形
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x66, 0x66, 0x66)
    line.line.fill.background()  # 无边框


def draw_polyline(slide, x1, y1, x2, y2):
    """绘制折线连接：从父节点底部向下，再水平，再向下到子节点顶部"""
    from pptx.enum.shapes import MSO_SHAPE

    # 计算中点y坐标
    mid_y = (y1 + y2) / 2

    # 绘制三条线段：垂直 -> 水平 -> 垂直
    # 第一段：从父节点底部向下
    line1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        x1 - Inches(0.01), y1,
        Inches(0.02), mid_y - y1
    )
    line1.fill.solid()
    line1.fill.fore_color.rgb = RGBColor(0x66, 0x66, 0x66)
    line1.line.fill.background()

    # 第二段：水平线
    line2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        min(x1, x2), mid_y - Inches(0.01),
        abs(x2 - x1), Inches(0.02)
    )
    line2.fill.solid()
    line2.fill.fore_color.rgb = RGBColor(0x66, 0x66, 0x66)
    line2.line.fill.background()

    # 第三段：向下到子节点
    line3 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        x2 - Inches(0.01), mid_y,
        Inches(0.02), y2 - mid_y
    )
    line3.fill.solid()
    line3.fill.fore_color.rgb = RGBColor(0x66, 0x66, 0x66)
    line3.line.fill.background()


def get_employee_stats_by_level3():
    """获取按三层部门、category、rank分组的人数统计"""
    conn = sqlite3.connect('employee.db')
    cursor = conn.cursor()
    cursor.execute('''
        SELECT department_level3, category, rank, COUNT(*) as count
        FROM employee
        WHERE department_level3 IS NOT NULL AND category IS NOT NULL AND rank IS NOT NULL
        GROUP BY department_level3, category, rank
        ORDER BY department_level3, category, rank
    ''')
    stats = cursor.fetchall()
    conn.close()
    return stats


def draw_employee_tables(slide, node_positions, tree_roots):
    """在幻灯片下半部分绘制员工统计表格"""
    # 获取员工统计数据
    stats = get_employee_stats_by_level3()

    # 构建统计字典: {dept_level3: {category: {rank: count}}}
    stats_dict = {}
    for dept, category, rank, count in stats:
        if dept not in stats_dict:
            stats_dict[dept] = {}
        if category not in stats_dict[dept]:
            stats_dict[dept][category] = {}
        stats_dict[dept][category][rank] = count

    # 获取所有三层部门节点（level=2），按x位置排序
    level3_nodes = [(pos['name'], pos['x_center'], pos['box_width'])
                    for pos in node_positions.values() if pos['level'] == 2]
    level3_nodes.sort(key=lambda x: x[1])  # 按x位置排序

    # 获取所有category（用于确定行数）
    all_categories = set()
    for dept_stats in stats_dict.values():
        all_categories.update(dept_stats.keys())
    categories = sorted(list(all_categories))

    # 配置参数
    slide_height = Inches(7.5)
    slide_width = Inches(13.333)

    # 表格样式参数
    cell_height = Inches(0.22)  # 每行高度
    table_height = cell_height * 4  # 标题 + 3个rank行
    rank_order = ['21', '20', '19']  # 从上到下的顺序（与数据库中的rank格式一致）

    # 计算行起始位置
    row_height = table_height + Inches(0.08)  # 每行高度（表格+间距）
    start_y = slide_height * 0.45  # 起始y位置

    # 绘制category标签（左侧）
    label_width = Inches(0.4)
    label_x = Inches(0.05)  # 更靠左

    # 为每个category绘制一行
    for cat_idx, category in enumerate(categories):
        row_y = start_y + cat_idx * row_height

        # 绘制category标签
        draw_category_label(slide, label_x, row_y, label_width, table_height, category)

        # 在该行内，按照三层部门的横向位置绘制表格
        for dept_name, x_center, box_width in level3_nodes:
            if dept_name not in stats_dict:
                continue
            if category not in stats_dict[dept_name]:
                continue

            # 表格x位置：与三层部门的x_center对齐
            # 表格宽度参考部门节点框宽度，但设置最小和最大值
            table_width = max(min(box_width * 0.9, Inches(1.2)), Inches(0.5))
            table_x = x_center - table_width / 2

            # 绘制表格
            draw_stats_table(slide, table_x, row_y, table_width, cell_height,
                           stats_dict[dept_name][category], rank_order)


def draw_category_label(slide, x, y, width, height, text):
    """绘制category标签（竖向纯文字）"""
    from pptx.enum.shapes import MSO_SHAPE

    # 每个字单独一行，实现竖向显示
    textbox = slide.shapes.add_textbox(x, y, width, height)
    tf = textbox.text_frame
    tf.word_wrap = False

    # 将文字竖向排列：每个字符作为单独段落
    for i, char in enumerate(text):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = char
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0)
        p.space_after = Pt(0)


def draw_stats_table(slide, x, y, width, cell_height, rank_data, rank_order):
    """绘制一个统计表格 (1列 x 3行: P21/P20/P19)"""
    from pptx.enum.shapes import MSO_SHAPE

    # 表格高度
    table_height = cell_height * 3  # 3个rank行

    # 绘制表格外边框
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, table_height)
    border.fill.solid()
    border.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # 白色背景
    border.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    border.line.width = Pt(0.5)

    # 绘制rank行
    for i, rank in enumerate(rank_order):
        row_y = y + i * cell_height

        # 绘制单元格
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, row_y, width, cell_height)
        cell.fill.solid()
        count = rank_data.get(rank, 0)
        if count > 0:
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)  # 浅蓝色背景
        else:
            cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)  # 浅灰色背景
        cell.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        cell.line.width = Pt(0.5)

        # 添加文字
        tf = cell.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = f"{rank}: {count}"
        p.font.size = Pt(7)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(3)


def generate_documentation(output_file, node_positions, level_colors, slide, level_widths):
    """生成PPT元素说明文档"""
    doc_path = Path(__file__).parent / 'PPT元素说明.md'

    # 统计连接线数量
    line_count = 0
    for shape in slide.shapes:
        if hasattr(shape, 'line') and shape.shape_type == 1:  # 矩形类型
            # 检查是否为线条（宽度或高度很小）
            if shape.width < Inches(0.1) or shape.height < Inches(0.1):
                line_count += 1

    # 按层级统计部门数量
    level_counts = {}
    for pos in node_positions.values():
        level = pos['level']
        level_counts[level] = level_counts.get(level, 0) + 1

    # 颜色映射
    level_color_names = ['蓝色', '绿色', '橙色', '红色', '紫色']

    # 计算动态宽度范围
    box_width_values = [pos['box_width'] for pos in node_positions.values()]
    min_width = min(box_width_values) if box_width_values else Inches(1.5)
    max_width = max(box_width_values) if box_width_values else Inches(1.5)

    # 获取员工统计信息
    stats = get_employee_stats_by_level3()
    stats_dict = {}
    for dept, category, rank, count in stats:
        if dept not in stats_dict:
            stats_dict[dept] = {}
        if category not in stats_dict[dept]:
            stats_dict[dept][category] = {}
        stats_dict[dept][category][rank] = count

    all_categories = set()
    for dept_stats in stats_dict.values():
        all_categories.update(dept_stats.keys())
    categories = sorted(list(all_categories))

    # 统计表格数量
    table_count = 0
    for dept_stats in stats_dict.values():
        table_count += len(dept_stats)

    content = f"""# PPT元素说明文档

**生成时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
**文件名**: {Path(output_file).name}

---

## 1. 幻灯片设置

| 属性 | 值 |
|------|-----|
| 宽度 | 13.333英寸 |
| 高度 | 7.5英寸 |
| 比例 | 16:9 |
| 布局 | 空白布局 |

---

## 2. 部门节点框（矩形）

**总数**: {len(node_positions)}个

**布局**: 部门树占幻灯片高度的40%，位于页面顶部

**宽度策略**: 根据同层部门数量动态调整，避免重叠

### 节点框样式

| 属性 | 值 |
|------|-----|
| 形状 | 矩形 |
| 宽度 | 动态计算 ({min_width:.2f} ~ {max_width:.2f} 英寸) |
| 高度 | 0.35英寸 |
| 边框颜色 | #333333 |
| 文字颜色 | 白色 #FFFFFF |
| 文字大小 | 动态调整 (7~11pt) |
| 文字样式 | 加粗 |
| 文字对齐 | 居中 |
| 文字截断 | 超过10字显示省略号 |

### 各层级统计

| 层级 | 颜色 | 色值 | 数量 |
|------|------|------|------|
"""

    for level in sorted(level_counts.keys()):
        color = level_colors[level]
        color_name = level_color_names[level] if level < len(level_color_names) else '未知'
        # RGBColor 使用 [0], [1], [2] 访问 RGB 分量
        color_hex = f"#{color[0]:02X}{color[1]:02X}{color[2]:02X}"
        content += f"| {level + 1}层 | {color_name} | {color_hex} | {level_counts[level]} |\n"

    content += f"""
---

## 3. 连接线

**形式**: 折线（父节点底部 → 向下 → 水平 → 向下 → 子节点顶部）

**线段总数**: 约{line_count}个（每条折线由3个细矩形组成）

### 连接线样式

| 属性 | 值 |
|------|-----|
| 颜色 | 深灰色 #666666 |
| 线宽 | 约0.02英寸 |
| 形式 | 三个细矩形拼接成折线 |

---

## 4. 员工统计表格（下半部分）

**布局**: 每个category一行，表格横向与三层部门对齐

**表格数量**: {table_count}个

**三层部门数量**: {len([p for p in node_positions.values() if p['level'] == 2])}个

**Category数量**: {len(categories)}个

### 表格样式

| 属性 | 值 |
|------|-----|
| 结构 | 1列 × 3行 |
| 行内容 | P21、P20、P19及人数 |
| 单元格高度 | 0.22英寸 |
| 边框颜色 | #AAAAAA |
| 有数据背景色 | 浅蓝色 #E8F4FD |
| 无数据背景色 | 浅灰色 #F8F8F8 |
| 字体大小 | 7pt |

### Category标签样式

| 属性 | 值 |
|------|-----|
| 位置 | 每行左侧 |
| 宽度 | 0.6英寸 |
| 背景色 | 蓝色 #4A90D9 |
| 文字颜色 | 白色 |
| 字体大小 | 10pt |

### Category列表

"""

    for cat in categories:
        content += f"- {cat}\n"

    content += f"""
---

## 5. 元素总计

| 元素类型 | 数量 |
|----------|------|
| 矩形（部门框） | {len(node_positions)}个 |
| 连接线段 | 约{line_count}个 |
| Category标签 | {len(categories)}个 |
| 员工统计表格 | {table_count}个 |

---

## 6. 部门详情

"""

    # 按层级列出部门
    for level in sorted(level_counts.keys()):
        level_name = level_color_names[level] if level < len(level_color_names) else '未知'
        content += f"### {level + 1}层部门（{level_name}）\n\n"
        for pos in node_positions.values():
            if pos['level'] == level:
                content += f"- {pos['name']}\n"
        content += "\n"

    # 写入文件
    with open(doc_path, 'w', encoding='utf-8') as f:
        f.write(content)

    print(f'元素说明文档已更新: {doc_path}')


def main():
    # 从数据库获取部门数据
    departments = get_departments_from_db()

    # 构建部门树
    tree_roots = build_department_tree(departments)

    # 创建PPT (16:9 比例)
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 宽度
    prs.slide_height = Inches(7.5)    # 16:9 高度

    # 绘制树状图并获取节点位置信息
    node_positions, level_colors, slide, level_widths = draw_tree_with_return(prs, tree_roots)

    # 绘制员工统计表格（在幻灯片下半部分）
    draw_employee_tables(slide, node_positions, tree_roots)

    # 创建history目录
    project_dir = Path(__file__).parent
    history_dir = project_dir / 'history'
    history_dir.mkdir(exist_ok=True)

    # 移动旧的organization*.pptx文件到history目录
    for old_file in project_dir.glob('organization*.pptx'):
        dest = history_dir / old_file.name
        if dest.exists():
            dest.unlink()  # 删除history中同名文件
        shutil.move(str(old_file), str(dest))
        print(f'已将旧文件移动到: {dest}')

    # 生成新文件名
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    output_file = project_dir / f'organization_{timestamp}.pptx'

    # 保存PPT
    prs.save(str(output_file))
    print(f'PPT已生成: {output_file}')

    # 生成元素说明文档
    generate_documentation(output_file, node_positions, level_colors, slide, level_widths)


if __name__ == '__main__':
    main()

"""
Excel数据生成PPT主程序
流程：Excel → SQLite → PPT
"""
import sys
import os

# 获取程序所在目录（打包后exe所在目录）
if getattr(sys, 'frozen', False):
    # 打包后的exe
    PROGRAM_DIR = os.path.dirname(sys.executable)
else:
    # 开发环境
    PROGRAM_DIR = os.path.dirname(os.path.abspath(__file__))

import sqlite3
import pandas as pd
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def excel_to_sqlite(excel_file, db_file):
    """将Excel数据导入SQLite"""
    print(f'正在读取Excel文件: {excel_file}')

    if not os.path.exists(excel_file):
        print(f'错误: Excel文件不存在 - {excel_file}')
        return False

    xls = pd.ExcelFile(excel_file)
    print(f'包含的Sheet: {xls.sheet_names}')

    conn = sqlite3.connect(db_file)
    cursor = conn.cursor()

    try:
        # 导入department表
        if 'department' in xls.sheet_names:
            dept_df = pd.read_excel(xls, sheet_name='department')
            print(f'导入department表，共{len(dept_df)}条记录')
            cursor.execute('DELETE FROM department')
            for _, row in dept_df.iterrows():
                cursor.execute('''
                    INSERT INTO department (id, parent_id, name, level)
                    VALUES (?, ?, ?, ?)
                ''', (
                    int(row['id']) if pd.notna(row['id']) else None,
                    int(row['parent_id']) if pd.notna(row['parent_id']) else None,
                    row['name'] if pd.notna(row['name']) else None,
                    int(row['level']) if pd.notna(row['level']) else None
                ))

        # 导入employee表
        if 'employee' in xls.sheet_names:
            emp_df = pd.read_excel(xls, sheet_name='employee')
            print(f'导入employee表，共{len(emp_df)}条记录')
            cursor.execute('DELETE FROM employee')
            for _, row in emp_df.iterrows():
                cursor.execute('''
                    INSERT INTO employee (id, name, employee_number, department_level1,
                                        department_level2, department_level3,
                                        department_level4, department_level5, rank, category)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                ''', (
                    int(row['id']) if pd.notna(row['id']) else None,
                    row['name'] if pd.notna(row['name']) else None,
                    row['employee_number'] if pd.notna(row['employee_number']) else None,
                    row['department_level1'] if pd.notna(row['department_level1']) else None,
                    row['department_level2'] if pd.notna(row['department_level2']) else None,
                    row['department_level3'] if pd.notna(row['department_level3']) else None,
                    row['department_level4'] if pd.notna(row['department_level4']) else None,
                    row['department_level5'] if pd.notna(row['department_level5']) else None,
                    row['rank'] if pd.notna(row['rank']) else None,
                    row['category'] if pd.notna(row['category']) else None
                ))

        conn.commit()
        print('数据导入完成')
        return True

    except Exception as e:
        conn.rollback()
        print(f'导入失败: {e}')
        return False
    finally:
        conn.close()


def get_departments_from_db(db_file):
    """从数据库获取所有部门数据"""
    conn = sqlite3.connect(db_file)
    cursor = conn.cursor()
    cursor.execute('SELECT id, parent_id, name, level FROM department ORDER BY level, id')
    departments = cursor.fetchall()
    conn.close()
    return departments


def build_department_tree(departments):
    """构建部门树结构"""
    nodes = {d[0]: {'id': d[0], 'parent_id': d[1], 'name': d[2], 'level': d[3], 'children': []}
             for d in departments}
    roots = []
    for d in departments:
        node = nodes[d[0]]
        if d[1] is None:
            roots.append(node)
        else:
            parent = nodes.get(d[1])
            if parent:
                parent['children'].append(node)
    return roots


def get_level_widths(nodes):
    """计算每层的最大宽度"""
    widths = {}
    def traverse(node, level):
        widths[level] = widths.get(level, 0) + 1
        for child in node['children']:
            traverse(child, level + 1)
    for root in nodes:
        traverse(root, 0)
    return widths


def draw_node(slide, x, y, width, height, text, fill_color):
    """绘制一个节点"""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)

    if len(text) > 10:
        display_text = text[:10] + '...'
    else:
        display_text = text

    width_inches = width / 914400
    char_count = len(display_text)
    available_width = width_inches - 0.1

    if char_count > 0:
        max_char_width = available_width / char_count
        font_size = min(11, max(7, int(max_char_width / 0.014)))
    else:
        font_size = 11

    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = display_text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(6)


def draw_polyline(slide, x1, y1, x2, y2):
    """绘制折线连接"""
    from pptx.enum.shapes import MSO_SHAPE
    mid_y = (y1 + y2) / 2

    line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x1 - Inches(0.01), y1, Inches(0.02), mid_y - y1)
    line1.fill.solid()
    line1.fill.fore_color.rgb = RGBColor(0x66, 0x66, 0x66)
    line1.line.fill.background()

    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, min(x1, x2), mid_y - Inches(0.01), abs(x2 - x1), Inches(0.02))
    line2.fill.solid()
    line2.fill.fore_color.rgb = RGBColor(0x66, 0x66, 0x66)
    line2.line.fill.background()

    line3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x2 - Inches(0.01), mid_y, Inches(0.02), y2 - mid_y)
    line3.fill.solid()
    line3.fill.fore_color.rgb = RGBColor(0x66, 0x66, 0x66)
    line3.line.fill.background()


def draw_tree_with_return(presentation, nodes):
    """在PPT上绘制树状图，并返回节点位置信息"""
    from pptx.enum.shapes import MSO_SHAPE

    blank_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_layout)

    slide_height = Inches(7.5)
    slide_width = Inches(13.333)
    tree_height = slide_height * 0.4
    box_height = Inches(0.35)
    margin = Inches(0.3)

    level_widths = get_level_widths(nodes)
    max_level = max(level_widths.keys()) if level_widths else 0

    available_width = slide_width - 2 * margin
    max_count = max(level_widths.values()) if level_widths else 1
    max_box_width = (available_width - max_count * Inches(0.15)) / max_count
    max_box_width = min(max_box_width, Inches(2.0))

    level_height = tree_height / (max_level + 1)
    start_y = Inches(0.2)

    level_colors = [
        RGBColor(0x4A, 0x90, 0xD9),
        RGBColor(0x5C, 0xB8, 0x5C),
        RGBColor(0xF0, 0xAD, 0x4E),
        RGBColor(0xD9, 0x5F, 0x5F),
        RGBColor(0x9B, 0x59, 0xB6),
    ]

    node_positions = {}
    box_widths = {}

    def calculate_positions(node, depth, x_start, x_end):
        y = start_y + depth * level_height
        space = x_end - x_start
        gap = Inches(0.1)
        box_w = min(space - gap, max_box_width)
        box_w = max(box_w, Inches(0.6))

        x_center = (x_start + x_end) / 2
        x = x_center - box_w / 2

        node_positions[node['id']] = {
            'x': x,
            'y': y,
            'x_center': x_center,
            'y_center': y + box_height / 2,
            'box_right': x + box_w,
            'box_bottom': y + box_height,
            'name': node['name'],
            'level': depth,
            'box_width': box_w
        }
        box_widths[node['id']] = box_w

        if node['children']:
            num_children = len(node['children'])
            child_width = (x_end - x_start) / num_children
            for i, child in enumerate(node['children']):
                child_x_start = x_start + i * child_width
                child_x_end = child_x_start + child_width
                calculate_positions(child, depth + 1, child_x_start, child_x_end)

    for root in nodes:
        calculate_positions(root, 0, margin, slide_width - margin)

    def draw_tree_recursive(node):
        pos = node_positions[node['id']]
        for child in node['children']:
            child_pos = node_positions[child['id']]
            draw_polyline(slide, pos['x_center'], pos['box_bottom'], child_pos['x_center'], child_pos['y'])
            draw_tree_recursive(child)

    for root in nodes:
        draw_tree_recursive(root)

    for node_id, pos in node_positions.items():
        draw_node(slide, pos['x'], pos['y'], box_widths[node_id], box_height,
                 pos['name'], level_colors[pos['level']])

    return node_positions, slide


def get_employee_stats_by_level3(db_file):
    """获取按三层部门、category、rank分组的人数统计"""
    conn = sqlite3.connect(db_file)
    cursor = conn.cursor()
    cursor.execute('''
        SELECT department_level3, category, rank, COUNT(*) as count
        FROM employee
        WHERE department_level3 IS NOT NULL AND category IS NOT NULL AND rank IS NOT NULL
        GROUP BY department_level3, category, rank
        ORDER BY department_level3, category, rank
    ''')
    stats = cursor.fetchall()
    conn.close()
    return stats


def draw_category_label(slide, x, y, width, height, text):
    """绘制category标签（竖向纯文字）"""
    textbox = slide.shapes.add_textbox(x, y, width, height)
    tf = textbox.text_frame
    tf.word_wrap = False

    for i, char in enumerate(text):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = char
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0)
        p.space_after = Pt(0)


def draw_stats_table(slide, x, y, width, cell_height, rank_data, rank_order):
    """绘制一个统计表格"""
    from pptx.enum.shapes import MSO_SHAPE

    table_height = cell_height * 3

    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, table_height)
    border.fill.solid()
    border.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    border.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    border.line.width = Pt(0.5)

    for i, rank in enumerate(rank_order):
        row_y = y + i * cell_height
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, row_y, width, cell_height)
        cell.fill.solid()
        count = rank_data.get(rank, 0)
        if count > 0:
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
        else:
            cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
        cell.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        cell.line.width = Pt(0.5)

        tf = cell.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = f"{rank}: {count}"
        p.font.size = Pt(7)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(3)


def draw_employee_tables(slide, node_positions, db_file):
    """在幻灯片下半部分绘制员工统计表格"""
    stats = get_employee_stats_by_level3(db_file)

    stats_dict = {}
    for dept, category, rank, count in stats:
        if dept not in stats_dict:
            stats_dict[dept] = {}
        if category not in stats_dict[dept]:
            stats_dict[dept][category] = {}
        stats_dict[dept][category][rank] = count

    level3_nodes = [(pos['name'], pos['x_center'], pos['box_width'])
                    for pos in node_positions.values() if pos['level'] == 2]
    level3_nodes.sort(key=lambda x: x[1])

    all_categories = set()
    for dept_stats in stats_dict.values():
        all_categories.update(dept_stats.keys())
    categories = sorted(list(all_categories))

    slide_height = Inches(7.5)

    cell_height = Inches(0.22)
    table_height = cell_height * 4
    rank_order = ['21', '20', '19']

    row_height = table_height + Inches(0.08)
    start_y = slide_height * 0.45

    label_width = Inches(0.4)
    label_x = Inches(0.05)

    for cat_idx, category in enumerate(categories):
        row_y = start_y + cat_idx * row_height
        draw_category_label(slide, label_x, row_y, label_width, table_height, category)

        for dept_name, x_center, box_width in level3_nodes:
            if dept_name not in stats_dict:
                continue
            if category not in stats_dict[dept_name]:
                continue

            table_width = max(min(box_width * 0.9, Inches(1.2)), Inches(0.5))
            table_x = x_center - table_width / 2

            draw_stats_table(slide, table_x, row_y, table_width, cell_height,
                           stats_dict[dept_name][category], rank_order)


def generate_ppt(db_file, output_dir):
    """生成PPT"""
    print('正在生成PPT...')

    departments = get_departments_from_db(db_file)
    tree_roots = build_department_tree(departments)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    node_positions, slide = draw_tree_with_return(prs, tree_roots)
    draw_employee_tables(slide, node_positions, db_file)

    # 创建history目录
    history_dir = Path(output_dir) / 'history'
    history_dir.mkdir(exist_ok=True)

    # 移动旧的PPT文件
    for old_file in Path(output_dir).glob('organization*.pptx'):
        dest = history_dir / old_file.name
        if dest.exists():
            dest.unlink()
        import shutil
        shutil.move(str(old_file), str(dest))
        print(f'已将旧文件移动到: {dest}')

    # 生成新文件
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    output_file = Path(output_dir) / f'organization_{timestamp}.pptx'

    prs.save(str(output_file))
    print(f'PPT已生成: {output_file}')

    return output_file


def init_database(db_file):
    """初始化数据库表结构"""
    conn = sqlite3.connect(db_file)
    cursor = conn.cursor()

    cursor.execute('''
        CREATE TABLE IF NOT EXISTS department (
            id INTEGER PRIMARY KEY,
            parent_id INTEGER,
            name TEXT,
            level INTEGER
        )
    ''')

    cursor.execute('''
        CREATE TABLE IF NOT EXISTS employee (
            id INTEGER PRIMARY KEY,
            name TEXT NOT NULL,
            employee_number TEXT,
            department_level1 TEXT,
            department_level2 TEXT,
            department_level3 TEXT,
            department_level4 TEXT,
            department_level5 TEXT,
            rank TEXT,
            category TEXT
        )
    ''')

    conn.commit()
    conn.close()


def main():
    print('=' * 50)
    print('Excel数据生成PPT工具')
    print('=' * 50)
    print()

    # 文件路径
    excel_file = os.path.join(PROGRAM_DIR, 'data_template.xlsx')
    db_file = os.path.join(PROGRAM_DIR, 'employee.db')

    # 初始化数据库
    init_database(db_file)

    # Step 1: Excel → SQLite
    print('【Step 1】Excel数据导入SQLite')
    if not excel_to_sqlite(excel_file, db_file):
        print('数据导入失败，程序退出')
        return

    print()

    # Step 2: SQLite → PPT
    print('【Step 2】生成PPT')
    output_file = generate_ppt(db_file, PROGRAM_DIR)

    print()
    print('=' * 50)
    print('处理完成!')
    print(f'PPT文件: {output_file}')
    print('=' * 50)


if __name__ == '__main__':
    main()

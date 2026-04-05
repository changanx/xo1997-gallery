"""
PPT 生成器
"""
from typing import Tuple, Dict, List
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


class PPTGenerator:
    """组织架构 PPT 生成器"""

    # 层级颜色
    LEVEL_COLORS = [
        RGBColor(0x4A, 0x90, 0xD9),  # 蓝色
        RGBColor(0x5C, 0xB8, 0x5C),  # 绿色
        RGBColor(0xF0, 0xAD, 0x4E),  # 橙色
        RGBColor(0xD9, 0x5F, 0x5F),  # 红色
        RGBColor(0x9B, 0x59, 0xB6),  # 紫色
    ]

    def __init__(self):
        self.prs = None

    def generate(self, tree: List[dict], stats: List[Tuple], output_path: str) -> Tuple[bool, str]:
        """
        生成 PPT

        Args:
            tree: 部门树结构
            stats: 员工统计数据
            output_path: 输出文件路径

        Returns:
            (成功与否, 消息)
        """
        try:
            self.prs = Presentation()
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)

            # 创建空白幻灯片
            blank_layout = self.prs.slide_layouts[6]
            slide = self.prs.slides.add_slide(blank_layout)

            # 计算节点位置
            node_positions = self._calculate_positions(tree)

            # 绘制部门树
            self._draw_tree(slide, tree, node_positions)

            # 绘制员工统计表格
            self._draw_stats_tables(slide, node_positions, stats)

            # 保存文件
            self.prs.save(output_path)
            return True, f"PPT 已生成: {output_path}"

        except Exception as e:
            return False, f"生成失败: {str(e)}"

    def _calculate_positions(self, nodes: List[dict]) -> Dict[int, dict]:
        """计算节点位置"""
        positions = {}
        slide_width = Inches(13.333)
        slide_height = Inches(7.5)
        tree_height = slide_height * 0.4
        box_height = Inches(0.35)
        margin = Inches(0.3)

        # 计算每层宽度
        level_widths = self._get_level_widths(nodes)
        max_level = max(level_widths.keys()) if level_widths else 0

        available_width = slide_width - 2 * margin
        max_count = max(level_widths.values()) if level_widths else 1
        max_box_width = min((available_width - max_count * Inches(0.15)) / max_count, Inches(2.0))

        level_height = tree_height / (max_level + 1)
        start_y = Inches(0.2)

        def calc(node, depth, x_start, x_end):
            y = start_y + depth * level_height
            space = x_end - x_start
            gap = Inches(0.1)
            box_w = max(min(space - gap, max_box_width), Inches(0.6))

            x_center = (x_start + x_end) / 2
            x = x_center - box_w / 2

            positions[node['id']] = {
                'x': x,
                'y': y,
                'x_center': x_center,
                'y_center': y + box_height / 2,
                'box_bottom': y + box_height,
                'name': node['name'],
                'level': depth,
                'box_width': box_w,
            }

            if node['children']:
                num_children = len(node['children'])
                child_width = (x_end - x_start) / num_children
                for i, child in enumerate(node['children']):
                    calc(child, depth + 1, x_start + i * child_width, x_start + (i + 1) * child_width)

        for root in nodes:
            calc(root, 0, margin, slide_width - margin)

        return positions

    def _get_level_widths(self, nodes: List[dict]) -> Dict[int, int]:
        """计算每层的最大宽度"""
        widths = {}

        def traverse(node, level):
            widths[level] = widths.get(level, 0) + 1
            for child in node['children']:
                traverse(child, level + 1)

        for root in nodes:
            traverse(root, 0)
        return widths

    def _draw_tree(self, slide, nodes: List[dict], positions: Dict[int, dict]):
        """绘制部门树"""
        box_height = Inches(0.35)

        def draw_connections(node):
            pos = positions[node['id']]
            for child in node['children']:
                child_pos = positions[child['id']]
                self._draw_polyline(slide, pos['x_center'], pos['box_bottom'],
                                   child_pos['x_center'], child_pos['y'])
                draw_connections(child)

        for root in nodes:
            draw_connections(root)

        # 绘制节点
        for node_id, pos in positions.items():
            self._draw_node(slide, pos['x'], pos['y'], pos['box_width'], box_height,
                           pos['name'], self.LEVEL_COLORS[pos['level'] % len(self.LEVEL_COLORS)])

    def _draw_node(self, slide, x, y, width, height, text, fill_color):
        """绘制节点"""
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)

        # 截断文本
        display_text = text[:10] + '...' if len(text) > 10 else text

        # 计算字体大小
        width_inches = width / 914400
        char_count = len(display_text)
        if char_count > 0:
            max_char_width = (width_inches - 0.1) / char_count
            font_size = min(11, max(7, int(max_char_width / 0.014)))
        else:
            font_size = 11

        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = display_text
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)

    def _draw_polyline(self, slide, x1, y1, x2, y2):
        """绘制折线连接"""
        mid_y = (y1 + y2) / 2

        # 垂直线1
        line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x1 - Inches(0.01), y1, Inches(0.02), mid_y - y1)
        line1.fill.solid()
        line1.fill.fore_color.rgb = RGBColor(0x66, 0x66, 0x66)
        line1.line.fill.background()

        # 水平线
        line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, min(x1, x2), mid_y - Inches(0.01), abs(x2 - x1), Inches(0.02))
        line2.fill.solid()
        line2.fill.fore_color.rgb = RGBColor(0x66, 0x66, 0x66)
        line2.line.fill.background()

        # 垂直线2
        line3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x2 - Inches(0.01), mid_y, Inches(0.02), y2 - mid_y)
        line3.fill.solid()
        line3.fill.fore_color.rgb = RGBColor(0x66, 0x66, 0x66)
        line3.line.fill.background()

    def _draw_stats_tables(self, slide, node_positions: Dict, stats: List[Tuple]):
        """绘制员工统计表格"""
        if not stats:
            return

        # 整理统计数据
        stats_dict = {}
        categories = set()
        for dept, category, rank, count in stats:
            if dept not in stats_dict:
                stats_dict[dept] = {}
            if category not in stats_dict[dept]:
                stats_dict[dept][category] = {}
            stats_dict[dept][category][rank] = count
            categories.add(category)

        categories = sorted(categories)

        # 获取第三层部门节点
        level3_nodes = [(pos['name'], pos['x_center'], pos['box_width'])
                        for pos in node_positions.values() if pos['level'] == 2]
        level3_nodes.sort(key=lambda x: x[1])

        slide_height = Inches(7.5)
        cell_height = Inches(0.22)
        table_height = cell_height * 3
        rank_order = ['21', '20', '19']

        row_height = table_height + Inches(0.08)
        start_y = slide_height * 0.45

        for cat_idx, category in enumerate(categories):
            row_y = start_y + cat_idx * row_height

            # 绘制类别标签
            self._draw_category_label(slide, Inches(0.05), row_y, Inches(0.4), table_height, category)

            for dept_name, x_center, box_width in level3_nodes:
                if dept_name not in stats_dict or category not in stats_dict[dept_name]:
                    continue

                table_width = max(min(box_width * 0.9, Inches(1.2)), Inches(0.5))
                table_x = x_center - table_width / 2

                self._draw_stats_table(slide, table_x, row_y, table_width, cell_height,
                                       stats_dict[dept_name][category], rank_order)

    def _draw_category_label(self, slide, x, y, width, height, text):
        """绘制类别标签"""
        textbox = slide.shapes.add_textbox(x, y, width, height)
        tf = textbox.text_frame
        tf.word_wrap = False

        for i, char in enumerate(text):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = char
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

    def _draw_stats_table(self, slide, x, y, width, cell_height, rank_data, rank_order):
        """绘制统计表格"""
        table_height = cell_height * 3

        # 外框
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, table_height)
        border.fill.solid()
        border.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        border.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        border.line.width = Pt(0.5)

        for i, rank in enumerate(rank_order):
            row_y = y + i * cell_height
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, row_y, width, cell_height)
            cell.fill.solid()
            count = rank_data.get(rank, 0)
            if count > 0:
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
            else:
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
            cell.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            cell.line.width = Pt(0.5)

            tf = cell.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = f"{rank}: {count}"
            p.font.size = Pt(7)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(3)
